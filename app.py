import streamlit as st
from anthropic import Anthropic
from dotenv import load_dotenv
import base64
import io
from pptx import Presentation
from docx import Document

load_dotenv()

client = Anthropic()

SYSTEM_PROMPT = """
You are a Socratic science tutor for high school students (grades 9-12). Your role is to help students genuinely understand science — not just memorize answers — by guiding them through questions rather than explanations.

## Your Core Philosophy
You believe that a student who arrives at an answer through their own thinking understands it far more deeply than one who was just told the answer. Your job is to be the guide, not the answer key.

## How You Behave
- ALWAYS respond to a student's question with a question of your own that nudges them toward the answer
- Start by probing what they already know — never assume they know nothing, and never assume they know everything
- Ask one question at a time. Don't overwhelm them.
- When a student gives a partially correct answer, acknowledge what's right, then probe the part that's incomplete or wrong
- When a student gives a wrong answer, don't say "wrong" — instead ask a question that creates a contradiction with their reasoning so they discover the error themselves
- Use real-world analogies and examples to make abstract science concepts feel tangible
- Keep your tone warm, encouraging, and patient — like a favorite teacher who believes in them

## Your Stance on Giving Answers
- Never give the answer outright, even if the student asks you directly
- If the student has been genuinely struggling through 3-4 exchanges and is clearly stuck, you may offer a targeted hint — not the answer, but a clue that unlocks the next step
- A good hint points them toward the right question to ask themselves, not the answer itself
- If a student is completely lost, break the problem into a smaller piece and start there

## Subjects You Cover
Biology, Chemistry, Physics, Earth Science — all standard high school science curriculum.

## What You Never Do
- Never write out full explanations unprompted
- Never give multiple-choice options as a shortcut
- Never make a student feel stupid for not knowing something
- Never go off-topic into non-science areas
- Never summarize what was learned at the end of a session — that is the student's job

## How You Check In During a Session
Every 3-4 exchanges, pause the Socratic questioning and ask the student to synthesize what they've figured out so far before moving on. Use natural language like: "Before we go further — what would you say you've worked out so far? Put it in your own words." This keeps them from losing track as the conversation builds. After they respond, affirm what's right, gently correct gaps, then continue.

## How You End a Session
When a conversation feels like it's wrapping up, or when a student says they're done or says thank you, do NOT summarize what was covered. Instead, turn it back to them with something like: "Before you go — in your own words, can you tell me the most important thing you figured out today? Try to explain it like you're teaching it to a friend." Only after they give their summary should you affirm what they got right and gently fill in anything they missed.
"""

# Page config
st.set_page_config(
    page_title="Socratic Science Tutor",
    page_icon="🔬",
    layout="centered"
)

# Header
st.title("🔬 Socratic Science Tutor")
st.caption("I won't give you the answer — but I'll help you find it yourself.")

# Initialize chat history
if "messages" not in st.session_state:
    st.session_state.messages = []
    # Get opening message from tutor
    response = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=1024,
        system=SYSTEM_PROMPT,
        messages=[{"role": "user", "content": "Hello, I just started a session."}]
    )
    opening = response.content[0].text
    st.session_state.messages.append({"role": "assistant", "content": opening})

# Display chat history
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# File uploader
uploaded_file = st.file_uploader(
    "Upload a file (optional — image, PDF, PowerPoint, Word doc...)",
    type=["png", "jpg", "jpeg", "pdf", "pptx", "docx"],
    label_visibility="visible"
)

# Chat input
if prompt := st.chat_input("Ask a science question..."):

    user_content = None

    if uploaded_file:
        file_type = uploaded_file.type

        # IMAGE
        if file_type in ["image/png", "image/jpeg"]:
            image_data = base64.standard_b64encode(uploaded_file.read()).decode("utf-8")
            user_content = [
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": file_type,
                        "data": image_data
                    }
                },
                {"type": "text", "text": prompt}
            ]
            with st.chat_message("user"):
                st.image(uploaded_file)
                st.markdown(prompt)

        # PDF
        elif file_type == "application/pdf":
            pdf_bytes = uploaded_file.read()
            pdf_data = base64.standard_b64encode(pdf_bytes).decode("utf-8")
            user_content = [
                {
                    "type": "document",
                    "source": {
                        "type": "base64",
                        "media_type": "application/pdf",
                        "data": pdf_data
                    }
                },
                {"type": "text", "text": prompt}
            ]
            with st.chat_message("user"):
                st.markdown(f"📄 *{uploaded_file.name}*")
                st.markdown(prompt)

        # PPTX
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(io.BytesIO(uploaded_file.read()))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                if slide_content:
                    slides_text.append(f"Slide {i+1}: " + " | ".join(slide_content))
            extracted = "\n".join(slides_text)
            user_content = [
                {"type": "text", "text": f"Here is the content of a PowerPoint presentation:\n\n{extracted}\n\n{prompt}"}
            ]
            with st.chat_message("user"):
                st.markdown(f"📊 *{uploaded_file.name}*")
                st.markdown(prompt)

        # DOCX
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(io.BytesIO(uploaded_file.read()))
            extracted = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            user_content = [
                {"type": "text", "text": f"Here is the content of a Word document:\n\n{extracted}\n\n{prompt}"}
            ]
            with st.chat_message("user"):
                st.markdown(f"📝 *{uploaded_file.name}*")
                st.markdown(prompt)

    # No file — plain text
    if user_content is None:
        user_content = prompt
        with st.chat_message("user"):
            st.markdown(prompt)

    st.session_state.messages.append({"role": "user", "content": user_content})

    # Get tutor response
    with st.chat_message("assistant"):
        with st.spinner("Thinking..."):
            response = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=1024,
                system=SYSTEM_PROMPT,
                messages=st.session_state.messages
            )
            reply = response.content[0].text
            st.markdown(reply)

    st.session_state.messages.append({"role": "assistant", "content": reply})